import base64
from email.mime.text import MIMEText
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from google.auth.transport.requests import Request
import pickle
import google.generativeai as genai
import json
import os
from dotenv import load_dotenv
from datetime import datetime, timezone, timedelta
import csv
import io
from email.mime.multipart import MIMEMultipart
from email.mime.base import MIMEBase
from email import encoders
import requests
from google.oauth2.credentials import Credentials
import uuid
import base64
from email.mime.text import MIMEText
# --- FIX: Load environment variables FIRST ---
load_dotenv()
SEARCH_CACHE = {}
# Now, read the keys from the loaded environment
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
GOOGLE_CSE_ID = os.getenv("GOOGLE_CSE_ID")
BROWSERLESS_API_KEY = os.getenv("BROWSERLESS_API_KEY")

# Configure the Gemini API
genai.configure(api_key=GEMINI_API_KEY)

SCOPES = ["https://www.googleapis.com/auth/gmail.send", "https://www.googleapis.com/auth/calendar.events"]

try:
    from pypdf import PdfReader
except ImportError:
    # Handle case where pypdf is not installed
    print("Warning: pypdf is not installed. PDF processing will not be available.")
    PdfReader = None

try:
    import docx
except ImportError:
    print("Warning: python-docx is not installed. DOCX processing will not be available.")
    docx = None

try:
    import pptx
except ImportError:
    print("Warning: python-pptx is not installed. PPTX processing will not be available.")
    pptx = None 
# Add this entire new function to agent.py
# Replace the existing extract_text_from_file function in agent.py (around line 38)

# In agent.py, replace the existing extract_text_from_file function with this:

def extract_text_from_file(uploaded_file):
    """Extracts text content from a file object (works with Streamlit and Flask)."""
    try:
        # Get the file type in a way that works for both frameworks
        file_type = getattr(uploaded_file, 'mimetype', getattr(uploaded_file, 'type', ''))
        content = ""

        # Reset file pointer to the beginning
        uploaded_file.seek(0)

        if file_type == "text/plain":
            content = uploaded_file.read().decode("utf-8")
        elif file_type == "application/pdf" and PdfReader:
            pdf_reader = PdfReader(uploaded_file)
            for page in pdf_reader.pages:
                content += page.extract_text() or ""
        elif file_type == "text/csv":
            content = uploaded_file.read().decode("utf-8")
        elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document" and docx:
            doc = docx.Document(uploaded_file)
            for para in doc.paragraphs:
                content += para.text + "\n"
        elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation" and pptx:
            pres = pptx.Presentation(uploaded_file)
            for slide in pres.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"
        else:
            return f"Unsupported file type: {file_type}. Or the required library is not installed."
        
        # Reset file pointer again in case it's used elsewhere
        uploaded_file.seek(0)
        return content
    except Exception as e:
        return f"Error reading file: {e}"
    
def web_search(query: str, num_results: int = 10):
    """
    Performs a Google web search after optimizing the user's query for better results.
    Caches results for 10 minutes.
    """
    # --- Caching Logic (from previous step, unchanged) ---
    current_time = datetime.now()
    normalized_query = query.strip().lower()
    if normalized_query in SEARCH_CACHE:
        cached_data = SEARCH_CACHE[normalized_query]
        if current_time - cached_data['timestamp'] < timedelta(minutes=10):
            print(f"INFO: Returning cached result for query: '{query}'")
            return cached_data['result']

    if not GOOGLE_API_KEY or not GOOGLE_CSE_ID:
        return "Cannot perform web search: The GOOGLE_API_KEY or GOOGLE_CSE_ID is not set in your environment."
        
    try:
        # --- NEW: Step 1 - Refine the user's query for better search results ---
        print(f"INFO: Original user query: '{query}'")
        refinement_model = genai.GenerativeModel("gemini-2.5-flash-lite")
        refinement_prompt = f"""
        You are a search query optimization expert. Your task is to convert a user's conversational request into a highly effective, keyword-based Google search query.
        For example:
        - User request: "tell me 10 national news of india"
        - Optimized query: "top 10 India national news headlines today"
        - User request: "what were the main points of the last G20 summit"
        - Optimized query: "G20 summit key outcomes summary"
        - User request: "who won the cricket world cup last year"
        - Optimized query: "ICC Cricket World Cup winner 2023"

        Do not include any explanations. Just provide the raw, optimized search query.

        User's request: "{query}"
        Optimized query: 
        """
        response = refinement_model.generate_content(refinement_prompt)
        optimized_query = response.text.strip()
        print(f"INFO: Optimized search query: '{optimized_query}'")

        # --- Step 2: Perform the Google Custom Search with the OPTIMIZED query ---
        service = build("customsearch", "v1", developerKey=GOOGLE_API_KEY)
        res = service.cse().list(q=optimized_query, cx=GOOGLE_CSE_ID, num=num_results).execute()
        
        items = res.get('items', [])
        if not items:
            return f"Web search for '{optimized_query}' returned no results. Please try rephrasing."

        raw_results_str = ""
        for item in items:
            raw_results_str += f"Title: {item.get('title', 'N/A')}\nSnippet: {item.get('snippet', 'N/A')}\n\n"

        # --- Step 3: Synthesize the answer (unchanged) ---
        synthesis_model = genai.GenerativeModel("gemini-2.5-flash-lite")
        synthesis_prompt = f"""
        You are a synthesis model. Your task is to answer the user's original query based *strictly* on the provided web search results.
        Adhere to any specific formatting or quantity requests in the original query (e.g., if the user asks for a 'list of 10', you MUST provide a list with 10 items).
        Do not include URLs or mention that the information came from a web search. Answer in a helpful, direct, and conversational tone.

        ---
        User's Original Query: "{query}"
        ---
        Web Search Results (from query "{optimized_query}"):
        {raw_results_str}
        ---
        Final Answer:
        """
        synthesis_response = synthesis_model.generate_content(synthesis_prompt)
        final_result = synthesis_response.text.strip()
        
        # --- Store result in cache (unchanged) ---
        SEARCH_CACHE[normalized_query] = {
            'timestamp': current_time,
            'result': final_result
        }
        
        return final_result

    except Exception as e:
        return f"An error occurred during the web search process: {e}"

def automate_browser(url: str, task: str):
    """
    Uses Browserless to perform an automated task on a webpage by generating
    and executing a Playwright script. Includes a self-healing mechanism for script errors.
    """
    if not BROWSERLESS_API_KEY:
        return "⚠️ Browserless API key is not set. Please add it to your .env file."

    print(f"INFO: Generating browser script for task: '{task}'")
    script_generation_model = genai.GenerativeModel("gemini-2.5-flash-lite")
    
    # --- PROMPT FOR THE INITIAL SCRIPT ATTEMPT ---
    script_prompt = f"""
    Based on the given URL and task, write a concise Playwright script to be run in Browserless.
    - The script MUST be a single JavaScript async arrow function `async ({{ page, context }}) => {{ ... }}`.
    - Use `await page.goto(context.url)` to navigate.
    - Use `await page.type(...)`, `await page.click(...)`, `await page.waitForSelector(...)`, etc., to interact with the page.
    - To handle multi-step tasks, you MUST chain `await` commands sequentially.
    - To return data, the script's final expression MUST be a return statement (e.g., `return await page.locator('...').innerText();`).
    - Do not include any explanations, comments, or markdown formatting—just the raw JavaScript code for the arrow function.

    ---
    **URL:** "{url}"
    **TASK:** "{task}"
    ---
    **YOUR SCRIPT (raw code only):**
    """
    
    try:
        # --- First Attempt ---
        initial_response = script_generation_model.generate_content(script_prompt)
        initial_script = initial_response.text.strip().replace("```javascript", "").replace("```", "")
        print(f"INFO: Generated script (Attempt 1):\n{initial_script}")

        api_url = f"https://chrome.browserless.io/function?token={BROWSERLESS_API_KEY}"
        payload = {"code": initial_script, "context": {"url": url}}
        
        response = requests.post(api_url, json=payload, timeout=60)
        response.raise_for_status() 
        
        return f"✅ Browser task completed successfully.\nResult: {response.text}"

    except requests.exceptions.HTTPError as e:
        # --- Self-Healing Logic for Script Errors (HTTP 400) ---
        if e.response and e.response.status_code == 400:
            error_from_browserless = e.response.text
            print(f"INFO: Script failed on first attempt. Attempting self-healing. Error: {error_from_browserless}")
            
            # --- PROMPT FOR THE SELF-HEALING ATTEMPT ---
            fix_prompt = f"""
            The previous Playwright script you generated failed. Analyze the original task, the faulty script, and the error message, then write a corrected script.

            **Original Task:** "{task}"
            
            **Faulty Script:**
            ```javascript
            {initial_script}
            ```
            
            **Error Message from Execution Engine:**
            "{error_from_browserless}"

            ---
            **INSTRUCTIONS:**
            - Read the error message carefully to understand why the script failed (e.g., selector not found, element is not visible, timeout).
            - Write a new, corrected JavaScript async arrow function `async ({{ page, context }}) => {{ ... }}`.
            - Do not include any explanations or comments, just the raw corrected script code.
            ---

            **Corrected Script:**
            """
            
            try:
                # --- Second (Self-Healing) Attempt ---
                print("INFO: Generating corrected script.")
                healing_response = script_generation_model.generate_content(fix_prompt)
                corrected_script = healing_response.text.strip().replace("```javascript", "").replace("```", "")
                print(f"INFO: Generated corrected script (Attempt 2):\n{corrected_script}")

                healing_payload = {"code": corrected_script, "context": {"url": url}}
                retry_response = requests.post(api_url, json=healing_payload, timeout=60)
                retry_response.raise_for_status()
                
                return f"✅ Browser task completed successfully after self-healing.\nResult: {retry_response.text}"
            
            except Exception as retry_e:
                return f"🚫 Browser automation failed even after a self-healing attempt. Final error: {retry_e}"
        
        # Handle other non-400 HTTP errors from the first attempt
        elif e.response and e.response.status_code == 403:
            return f"🚫 The browser automation service blocked the request to {url}. This often happens with high-security websites like Amazon or Google. Please try a different site."
        else:
            return f"An HTTP error occurred with the Browserless API request: {e}"
            
    except Exception as e:
        return f"An unexpected error occurred during browser automation: {e}"
    
def screenshot_website(url: str):
    """
    Uses Browserless to take a full-page screenshot of a URL and save it to the 'outputs' folder.
    """
    if not BROWSERLESS_API_KEY:
        return "⚠️ Browserless API key is not set. Please add it to your .env file."

    print(f"INFO: Taking screenshot of URL: '{url}'")
    try:
        # A specific Playwright script to take a screenshot and return it as a base64 string
        script = """
        async ({ page, context }) => {
            await page.goto(context.url, { waitUntil: 'networkidle2' });
            const buffer = await page.screenshot({ fullPage: true, type: 'png' });
            return buffer.toString('base64');
        }
        """

        api_url = f"https://chrome.browserless.io/function?token={BROWSERLESS_API_KEY}"
        payload = {"code": script, "context": {"url": url}}

        # Use a longer timeout as screenshots of large pages can take time
        response = requests.post(api_url, json=payload, timeout=90)
        response.raise_for_status()

        base64_image = response.text
        image_data = base64.b64decode(base64_image)

        # Ensure the output directory exists
        output_dir = 'outputs'
        os.makedirs(output_dir, exist_ok=True)

        # Create a unique filename
        filename = f"screenshot_{uuid.uuid4().hex[:8]}.png"
        filepath = os.path.join(output_dir, filename)

        with open(filepath, 'wb') as f:
            f.write(image_data)

        return f"✅ Screenshot saved successfully as '{filename}'. You can view it now."

    except requests.exceptions.HTTPError as e:
        if e.response.status_code == 403:
            return f"🚫 The browser automation service blocked the request to {url}. This often happens with high-security websites."
        return f"An HTTP error occurred with the Browserless API request: {e.response.text}"
    except Exception as e:
        return f"An unexpected error occurred during the screenshot process: {e}"
    
def get_google_service_from_token(token_info, service_name, service_version):
    """Creates a Google API service object from a user's token information."""
    creds = Credentials.from_authorized_user_info(token_info, SCOPES)
    service = build(service_name, service_version, credentials=creds)
    return service

# Replace the old function in agent.py with this one

# In agent.py, replace the create_calendar_event function with this:

def create_calendar_event(summary, start_time, end_time, description=None, recurrence=None, token_info=None, create_meet_link: bool = False):
    """Creates a calendar event, optionally with a Google Meet link."""
    try:
        service = get_google_service_from_token(token_info, 'calendar', 'v3')
        event = {
            'summary': summary,
            'description': description,
            'start': {'dateTime': start_time, 'timeZone': 'Asia/Kolkata'},
            'end': {'dateTime': end_time, 'timeZone': 'Asia/Kolkata'},
        }
        if recurrence:
            event['recurrence'] = [recurrence]
        
        if create_meet_link:
            conference_data = {
                "createRequest": {
                    "requestId": f"{uuid.uuid4().hex}",
                    "conferenceSolutionKey": {"type": "hangoutsMeet"}
                }
            }
            event['conferenceData'] = conference_data

        created_event = service.events().insert(calendarId='primary', body=event, conferenceDataVersion=1).execute()
        
        meet_link = created_event.get('hangoutLink', 'No meet link generated.')
        event_link = created_event.get('htmlLink')

        response = {
            "response": f"✅ Successfully created event: '{summary}'.\nEvent Link: {event_link}",
            "meet_link": meet_link
        }
        return response

    except Exception as e:
        # Return a dictionary for consistency, even on error
        return {"response": f"An error occurred while creating calendar event: {e}", "meet_link": None}
    

# In agent.py, REPLACE the entire get_upcoming_events function with this one:

def get_upcoming_events(max_results=10, time_min=None, time_max=None, token_info=None):
    """
    Fetches upcoming calendar events, optionally filtering by a time range.
    """
    try:
        service = get_google_service_from_token(token_info, 'calendar', 'v3')
        if time_min and not time_min.endswith('Z'):
            time_min += 'Z'
        if time_max and not time_max.endswith('Z'):
            time_max += 'Z'
            
        # If time_min is still not provided, default to the current time.
        if not time_min:
            time_min = datetime.now(timezone.utc).isoformat()
        events_result = service.events().list(
            calendarId='primary',
            timeMin=time_min,
            timeMax=time_max,
            maxResults=max_results,
            singleEvents=True,
            orderBy='startTime'
        ).execute()
        
        events = events_result.get('items', [])

        if not events:
            if time_max:
                 return "📅 You have no events scheduled for that day."
            return "📅 You have no upcoming events."

        response_text = "📅 Here are your events:\n"
        for event in events:
            start = event['start'].get('dateTime', event['start'].get('date'))
            try:
                # Improved parsing to handle different date formats from Google
                start_dt = datetime.fromisoformat(start.replace('Z', '+00:00'))
                start_formatted = start_dt.strftime('%a, %b %d at %I:%M %p')
            except ValueError:
                start_formatted = datetime.strptime(start, '%Y-%m-%d').strftime('%a, %b %d (All-day)')

            response_text += f"- **{event['summary']}** at {start_formatted}\n"
        
        return response_text
    except Exception as e:
        # --- UPGRADE: Provide a more detailed error message for easier debugging ---
        print(f"ERROR in get_upcoming_events: {e}")
        # Return a more specific error to the agent planner and user
        return f"An error occurred while fetching calendar events: The Google Calendar API rejected the request. Please ensure you are logged in with the correct account."
    
# In agent.py, replace the find_event_id, update_event, and delete_event functions

def find_event_id(event_summary, token_info=None):
    service = get_google_service_from_token(token_info, 'calendar', 'v3')
    now = datetime.now(timezone.utc)
    # Search within the next 30 days. This can be adjusted.
    time_max = now + timedelta(days=30)
    
    events_result = service.events().list(
        calendarId='primary', 
        q=event_summary,
        timeMin=now.isoformat(),
        timeMax=time_max.isoformat(),
        maxResults=5, 
        singleEvents=True,
        orderBy='startTime'
    ).execute()
    
    events = events_result.get('items', [])
    if not events:
        return None
    # Return the ID of the first match
    return events[0]['id']

def update_event(event_summary, new_start_time, new_end_time, token_info=None):
    """Finds an event by summary and updates its time."""
    try:
        service = get_google_service_from_token(token_info, 'calendar', 'v3')
        event_id = find_event_id(event_summary, token_info=token_info)
        
        if not event_id:
            return f"⚠️ Sorry, I couldn't find an event named '{event_summary}' to update."
        
        # First, get the existing event
        event = service.events().get(calendarId='primary', eventId=event_id).execute()
        
        # Update the start and end times
        event['start']['dateTime'] = new_start_time
        event['end']['dateTime'] = new_end_time
        
        updated_event = service.events().update(calendarId='primary', eventId=event_id, body=event).execute()
        return f"✅ Successfully updated event: '{updated_event.get('summary')}' is now at {new_start_time}."
    except Exception as e:
        return f"An error occurred while updating the event: {e}"

def delete_event(event_summary, token_info=None):
    """Finds an event by summary and deletes it."""
    try:
        service = get_google_service_from_token(token_info, 'calendar', 'v3')
        event_id = find_event_id(event_summary, token_info=token_info)
        
        if not event_id:
            return f"⚠️ Sorry, I couldn't find an event named '{event_summary}' to delete."
        
        service.events().delete(calendarId='primary', eventId=event_id).execute()
        return f"✅ Successfully deleted the event: '{event_summary}'."
    except Exception as e:
        return f"An error occurred while deleting the event: {e}"
    
    
# In agent.py, replace the existing send_email function with this:

def send_email(to, subject, body, attachment_object=None, token_info=None):
    """Sends an email, optionally with a file attachment."""
    try:
        service = get_google_service_from_token(token_info, 'gmail', 'v1')
        
        if attachment_object:
            message = MIMEMultipart()
            message["to"] = to
            message["subject"] = subject
            message.attach(MIMEText(body, "plain"))
            
            part = MIMEBase("application", "octet-stream")
            
            # Reset stream and read bytes (works for Flask/Werkzeug file objects)
            attachment_object.seek(0)
            file_bytes = attachment_object.read()
            part.set_payload(file_bytes)
            
            encoders.encode_base64(part)

            # Get filename in a way that works for both frameworks
            filename = getattr(attachment_object, 'filename', getattr(attachment_object, 'name', 'attachment'))
            part.add_header(
                "Content-Disposition",
                f"attachment; filename= {filename}",
            )
            message.attach(part)
            
        else:
            message = MIMEText(body)
            message["to"] = to
            message["subject"] = subject

        message["from"] = "me"
        encoded_message = base64.urlsafe_b64encode(message.as_bytes()).decode()
        create_message = {"raw": encoded_message}
        
        send_message = service.users().messages().send(userId="me", body=create_message).execute()
        return f"✅ Successfully sent email with ID: {send_message['id']}"
    except Exception as e:
        return f"An error occurred while sending email: {e}"
    
# In agent.py, DELETE your existing get_agent_plan function and REPLACE it with this:

def get_agent_plan(user_prompt, user_name, user_title, history=None, file_context=None, failure_context=None):
    model = genai.GenerativeModel("gemini-2.5-flash-lite")
    
    file_context_block = ""
    if file_context:
        file_context_block = f"""---
## File Context Provided
The user has uploaded a file. Prioritize using this context if the prompt refers to the file.
<FILE_CONTENT>
{file_context[:4000]}
</FILE_CONTENT>
---"""
        
    failure_context_block = ""
    if failure_context:
        failure_context_block = f"""---
## Previous Attempt Failed
Your last attempt to execute a plan failed. You MUST analyze the error and the original goal, then create a new, different plan to achieve the goal. Do not repeat the exact same failed step. For example, if browser automation failed, try using a web search instead.
<FAILURE_CONTEXT>
{failure_context}
</FAILURE_CONTEXT>
---"""
        
    history_block = ""
    if history:
        # --- THIS IS THE FIX ---
        # First, format the history into a single string
        formatted_history_list = []
        for message in history:
            role = "User" if message.get('role') == 'user' else "Assistant"
            content = message.get('content', '')
            formatted_history_list.append(f"{role}: {content}")
        history_string = "\n".join(formatted_history_list)
        
        # Then, insert that string into the f-string
        history_block = f"""---
## Conversation History
The following is the history of the current conversation. You MUST use it to understand context, answer follow-up questions, and maintain a coherent dialogue.
{history_string}
---"""


    current_time_iso = datetime.now(timezone.utc).isoformat()
    today = datetime.now()
    start_of_day_iso = today.replace(hour=0, minute=0, second=0, microsecond=0).isoformat() + "Z"
    end_of_day_iso = today.replace(hour=23, minute=59, second=59, microsecond=999999).isoformat() + "Z"

    tool_manual_str = """
- `analyze_file(task: str)`: Analyzes the content of the currently loaded file to perform a specific task (e.g., "summarize the document", "extract all email addresses", "what is the main topic of the text?"). This tool MUST be used if the user's query is about the file.
- `automate_browser(url: str, task: str)`: Loads a webpage and performs a specific task (e.g., "extract all the headings", "fill out the login form with username 'test' and password 'pass'"). Returns the result of the task.
- `screenshot_website(url: str)`: Takes a full-page screenshot of the provided URL and saves it as a PNG file.
- `web_search(query: str, num_results: int = 10)`: Searches the web for a given query. IMPORTANT: This tool can only return a maximum of 10 results due to API limitations. It returns a synthesized answer based on the top results. Use for current events, facts, or general knowledge.
- `send_email(to: str, subject: str, body: str, attach_file: bool = False)`: Sends an email. If `attach_file` is true and a file is in context, it will be attached. Use the user's signature if provided.
- `get_upcoming_events(max_results: int = 10, time_min: str = None, time_max: str = None)`: Fetches upcoming events from your Google Calendar.
- `Calendar(summary: str, start_time: str, end_time: str, description: str = None, recurrence: str = None, create_meet_link: bool = False)`: Creates a new event on your Google Calendar. Set `create_meet_link` to `True` to add a Google Meet.
- `update_event(event_summary: str, new_start_time: str, new_end_time: str)`: Finds an event by its summary and updates its start and end times.
- `delete_event(event_summary: str)`: Finds an event by its summary and deletes it.
- `conversational_response(response: str)`: Use this for simple conversational responses, greetings, or when no other tool is appropriate.
"""

    system_prompt = f"""
    You are a ReAct-style agent. Your goal is to break down a user's request into a sequence of tool calls and format this as a JSON array.

    ## Global Rules
1.  **Current Date & Time:** The current date is **{today.strftime('%A, %B %d, %Y')}**. The current UTC time is **{current_time_iso}**. You MUST use this information to accurately calculate dates and times for scheduling. For example, if today is Monday, then "tomorrow" is Tuesday. All generated `start_time` and `end_time` parameters MUST be in the full `YYYY-MM-DDTHH:MM:SS` ISO 8601 format.
2.  **JSON Array Output:** Your entire response MUST be a JSON array `[...]`.
3.  **Chaining Results:** Use `$STEP_X_RESULT` to refer to the output of a previous step. If the output is an object, you can access its properties using dot notation, e.g., `$STEP_1_RESULT.meet_link`.
    
    {history_block}
    {file_context_block}
    {failure_context_block}

    ---
    ## Tool Manual
    {tool_manual_str}
    ---
    ## Comprehensive Examples

    (Example 1: Multi-Step Browser Automation - unchanged)
    (Example 2: The Daily Briefing - unchanged)

    ### Example 3: Conversational Question
    User prompt: "Hi, what can you do for me?"
    Your response:
    [
      {{
        "tool": "conversational_response",
        "parameters": {{
            "response": "Hello! I am an AI assistant that can help you with a variety of tasks. I can browse the web, search for information, send emails, and manage your Google Calendar. I can also prepare a daily briefing for you with your events and the latest news. How can I help you today?"
        }}
      }}
    ]
    
    ### Example 4: Schedule a Google Meet and Email the Link
    User prompt: "schedule a 30 minute meeting with the marketing team for tomorrow at 2pm and send me the link"
    User Signature Provided: Name="Alex", Title="Project Manager"
    Your response:
    [
      {{
        "tool": "create_calendar_event",
        "parameters": {{
          "summary": "Meeting with the marketing team",
          "start_time": "2025-09-30T14:00:00",
          "end_time": "2025-09-30T14:30:00",
          "create_meet_link": true
        }}
      }},
      {{
        "tool": "send_email",
        "parameters": {{
          "to": "me",
          "subject": "Google Meet Link for 'Meeting with the marketing team'",
          "body": "Hi Alex,

Here is the Google Meet link you requested for your meeting:

$STEP_1_RESULT.meet_link

Best regards."
        }}
      }}
    ]
    """
    signature_block = f'User Signature Provided: Name="{user_name}", Title="{user_title}"' if user_name else ""
    full_prompt = f'{system_prompt}\n{signature_block}\nUser prompt: "{user_prompt}"\nYour response:'

    try:
        response = model.generate_content(full_prompt)
        plan_json = response.text.strip().replace("```json", "").replace("```", "")
        master_plan = json.loads(plan_json)
        return master_plan if isinstance(master_plan, list) else [master_plan]
    except Exception as e:
        print(f"Error generating plan: {e}")
        return [{"tool": "error", "parameters": {"response": "I had trouble creating a plan for that request. Please try rephrasing."}}]
    
    
# In agent.py, replace your entire run_agent_task_from_plan function with this:

def run_agent_task_from_plan(plan, file_context=None, file_object=None, token_info=None):
    if not plan or 'tool' not in plan:
        return {"tool": "error", "response": "🤔 Sorry, I couldn't determine a plan of action."}
    
    tool_name = plan['tool']
    parameters = plan.get('parameters', {})
    normalized = tool_name.strip().lower()
    
    def p(k, default=None):
        return parameters.get(k, default)

    if tool_name == "error":
        return {"tool": "error", "response": p('response')}
        
    response_text = ""
    tool_used = "unknown"

    if normalized in ('analyze_file', 'read_file'):
        tool_used = 'analyze_file'
        task = p('task', 'Provide a brief overview of the file.')
        if not file_context:
            response_text = "⚠️ You asked me to analyze a file, but no file has been provided or it could not be read."
        else:
            analysis_model = genai.GenerativeModel("gemini-2.5-flash-lite")
            analysis_prompt = f"""
            A user has provided a document for analysis.

            **User's Task:** "{task}"

            **Instructions for you, the AI:**
            1.  **Analyze Thoroughly:** Carefully read the document provided below to understand its content and structure.
            2.  **Fulfill the Task:** Complete the user's task based *strictly* on the information within the document.
            3.  **Structure Your Output:** Format your response for maximum clarity and readability. Use Markdown formatting:
                - Use headings (`##`) and subheadings (`###`) for different sections.
                - Use bullet points (`*`) or numbered lists (`1.`) for key items, summaries, or action points.
                - Use bold text (`**text**`) to highlight important terms or conclusions.
                - Write in well-structured paragraphs.
                - The final output should be clean, professional, and easy to digest.

            ---
            **DOCUMENT CONTENT:**
            {file_context}
            ---

            **Your Structured Response:**
            """
            response = analysis_model.generate_content(analysis_prompt)
            response_text = response.text
            
    elif normalized == 'automate_browser':
        tool_used = 'automate_browser'
        response_text = automate_browser(url=p('url'), task=p('task'))
    elif normalized == 'web_search':
        tool_used = 'web_search'
        response_text = web_search(query=p('query'))
    elif normalized == 'screenshot_website':
        tool_used = 'screenshot_website'
        response_text = screenshot_website(url=p('url'))    
    elif normalized in ('send_email', 'sendemail'):
        tool_used = 'send_email'
        should_attach = p('attach_file', False)
        attachment_to_send = None
        if should_attach:
            if file_object:
                attachment_to_send = file_object
            else:
                return {"tool": tool_used, "response": "⚠️ You asked to attach a file, but no file is currently loaded."}
        # --- THIS LINE IS UPDATED ---
        response_text = send_email(to=p('to'), subject=p('subject'), body=p('body'), attachment_object=attachment_to_send, token_info=token_info)
    elif normalized in ('create_calendar_event', 'calendar'):
        tool_used = 'create_calendar_event'
        # --- THIS LINE IS UPDATED ---
        response_text = create_calendar_event(summary=p('summary'), start_time=p('start_time'), end_time=p('end_time'), description=p('description'), recurrence=p('recurrence'), token_info=token_info, create_meet_link=p('create_meet_link', False))
    elif normalized in ('get_upcoming_events', 'events'):
        tool_used = 'get_upcoming_events'
        # --- THIS LINE IS UPDATED ---
        response_text = get_upcoming_events(max_results=p('max_results', 10), time_min=p('time_min'), time_max=p('time_max'), token_info=token_info)
    elif normalized in ('update_event', 'reschedule_event'):
        tool_used = 'update_event'
        # --- THIS LINE IS UPDATED ---
        response_text = update_event(event_summary=p('event_summary'), new_start_time=p('new_start_time'), new_end_time=p('new_end_time'), token_info=token_info)
    elif normalized in ('delete_event', 'remove_event'):
        tool_used = 'delete_event'
        # --- THIS LINE IS UPDATED ---
        response_text = delete_event(event_summary=p('event_summary'), token_info=token_info)
    elif normalized in ('conversational_response', 'conversation'):
        tool_used = 'conversational_response'
        response_text = p('response', "Hello! How can I help?")
    else:
        tool_used = 'error'
        response_text = f"⚠️ Unknown tool detected: {tool_name}"
        
    return {"tool": tool_used, "response": response_text}